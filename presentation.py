from os import listdir
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from scipy.ndimage import imread


def make_presentation(slides_dir_path, presentation_title):

	full_slide_paths = [slides_dir_path + slide_path for slide_path in sorted(listdir(slides_dir_path))]
	height, width, channels = imread(full_slide_paths[0]).shape
	height_inch = 0.010416666666819 * height
	width_inch = 0.010416666666819 * width

	presentation = Presentation()
	presentation.slide_height = Inches(height_inch)
	presentation.slide_width = Inches(width_inch)

	first_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
	first_slide.background.fill.solid()
	first_slide.background.fill.fore_color.rgb = RGBColor(224, 22, 120)
	preface = first_slide.shapes.add_textbox(left=Inches((width_inch - 3)/2), top=Inches((height_inch - 1.2)/2), height=Inches(1.2), width=Inches(3))
	title = preface.text_frame.add_paragraph()
	title.text = "Automation 101"
	title.alignment = PP_ALIGN.CENTER
	title.font.bold = True
	title.font.size = Pt(30)
	title.font.color.rgb = RGBColor(255, 255, 255)
	creds = preface.text_frame.add_paragraph()
	creds.text = "by bex on the beach"
	creds.alignment = PP_ALIGN.CENTER
	creds.font.color.rgb = RGBColor(209, 203, 207)

	for slide_path in full_slide_paths:

		curr_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
		curr_slide.shapes.add_picture(slide_path, left=Inches(0), top=Inches(0), height=Inches(height_inch), width=Inches(width_inch))

	presentation.save('{0}{1}.pptx'.format(slides_dir_path, presentation_title))
